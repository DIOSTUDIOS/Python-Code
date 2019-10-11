from pptx import Presentation
 
# 读取ppt
prs = Presentation('./志愿者慈善公益宣传PPT模板2.pptx')
 
# 查看一共几页
slides = prs.slides
number_pages = len(slides)
print(number_pages)
 
# 删除最后一页
rId = prs.slides._sldIdLst[-1].rId
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[-1]
 
# 保存新的ppt
prs.save('./志愿者慈善公益宣传PPT模板2.pptx')